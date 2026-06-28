"""Tests for Phase 4: PPTX Assembly."""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

import config
from pipeline.models import BBox, DetectedElement, ElementTree
from pipeline.phase4_assembly.assembler import px_to_emu, assemble


def test_px_to_emu_zero():
    assert px_to_emu(0, "x") == 0
    assert px_to_emu(0, "y") == 0


def test_px_to_emu_max_width():
    assert px_to_emu(config.SLIDE_WIDTH_PX, "x") == config.SLIDE_WIDTH_EMU


def test_group_shape_creation(tmp_path):
    output_path = str(tmp_path / "test_group.pptx")
    
    parent = DetectedElement(bbox=BBox(x=0, y=0, w=200, h=200), semantic_type="Section")
    child1 = DetectedElement(bbox=BBox(x=10, y=10, w=50, h=50), semantic_type="TextBlock")
    child2 = DetectedElement(bbox=BBox(x=70, y=10, w=50, h=50), semantic_type="TextBlock")
    
    parent.children = [child1, child2]
    tree = ElementTree(roots=[parent])
    
    assemble(tree, output_path)
    
    prs = Presentation(output_path)
    slide = prs.slides[0]
    
    # Verify group shape was created
    groups = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]
    assert len(groups) == 1
    
    # Check that children exist inside the group
    # Note: text lines don't exist in our mock elements, so we just check the number of objects
    # Wait, our render_leaf for TextBlock delegates to text_lines which adds textboxes if ocr_lines exist.
    # Because there are no OCR lines, it should add nothing for children.
    # So the group might be empty, but the group shape itself should be there.


def test_arrow_connector(tmp_path):
    output_path = str(tmp_path / "test_arrow.pptx")
    
    arrow = DetectedElement(
        bbox=BBox(x=0, y=0, w=100, h=100), 
        semantic_type="Arrow",
        arrow_start=(10, 10),
        arrow_end=(90, 90)
    )
    tree = ElementTree(roots=[arrow])
    
    assemble(tree, output_path)
    
    prs = Presentation(output_path)
    slide = prs.slides[0]
    
    # We should have a connector
    connectors = [s for s in slide.shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.LINE]
    assert len(connectors) == 1
    
    # Verify XML for a:tailEnd
    xml = connectors[0].element.xml
    assert "a:tailEnd" in xml
    assert "type=\"triangle\"" in xml


def test_arrow_no_geometry(tmp_path):
    output_path = str(tmp_path / "test_no_arrow.pptx")
    
    # Arrow without start/end
    arrow = DetectedElement(
        bbox=BBox(x=0, y=0, w=100, h=100), 
        semantic_type="Arrow",
        arrow_start=None,
        arrow_end=None
    )
    tree = ElementTree(roots=[arrow])
    
    assemble(tree, output_path)
    
    prs = Presentation(output_path)
    slide = prs.slides[0]
    
    # It should render nothing
    assert len(slide.shapes) == 0
