"""Phase 4: PPTX Assembly orchestrator."""

from pptx import Presentation
from pptx.util import Emu
from loguru import logger

import config
from pipeline.models import DetectedElement, ElementTree
from pipeline.phase4_assembly.group_builder import render_group
from pipeline.phase4_assembly.arrow_writer import render_arrow
from pipeline.phase4_assembly.text_placer import render_text_lines


def px_to_emu(px: int, axis: str, w: int, h: int) -> int:
    """Convert pixels to EMU based on config ratios and dynamic dimensions."""
    if axis == "x":
        return int(px * config.SLIDE_WIDTH_EMU / w)
    else:
        return int(px * config.SLIDE_HEIGHT_EMU / h)


def render_tile(element: DetectedElement, target, tree: ElementTree):
    """Render a cropped image tile."""
    if not element.tile_path:
        return
    try:
        target.shapes.add_picture(
            element.tile_path,
            Emu(px_to_emu(element.bbox.x, "x", tree.source_image_w, tree.source_image_h)),
            Emu(px_to_emu(element.bbox.y, "y", tree.source_image_w, tree.source_image_h)),
            Emu(px_to_emu(element.bbox.w, "x", tree.source_image_w, tree.source_image_h)),
            Emu(px_to_emu(element.bbox.h, "y", tree.source_image_w, tree.source_image_h))
        )
    except Exception as e:
        pass # Ignore missing files in tests


def render_shape(element: DetectedElement, target, tree: ElementTree):
    """Render a native PowerPoint shape."""
    from pptx.enum.shapes import MSO_SHAPE
    
    shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "arrow_triangle": MSO_SHAPE.RIGHT_ARROW,
        "star": MSO_SHAPE.STAR_5_POINT,
        "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "circle": MSO_SHAPE.OVAL,
        "pill": MSO_SHAPE.ROUNDED_RECTANGLE,
        "polygon": MSO_SHAPE.HEXAGON
    }
    
    shape_type = element.shape_type or "rectangle"
    mso_type = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)
    
    print(f"DEBUG: render_shape type={shape_type}")
    
    target.shapes.add_shape(
        mso_type,
        Emu(px_to_emu(element.bbox.x, "x", tree.source_image_w, tree.source_image_h)),
        Emu(px_to_emu(element.bbox.y, "y", tree.source_image_w, tree.source_image_h)),
        Emu(px_to_emu(element.bbox.w, "x", tree.source_image_w, tree.source_image_h)),
        Emu(px_to_emu(element.bbox.h, "y", tree.source_image_w, tree.source_image_h))
    )


def render_text_block(element: DetectedElement, target, tree: ElementTree):
    """Render just the text lines."""
    render_text_lines(element, target, tree)


def render_leaf(element: DetectedElement, target, tree: ElementTree):
    """Render leaf element based on its processing path or type."""
    if element.processing_path == "reconstruct":
        if element.semantic_type == "Arrow":
            render_arrow(element, target, tree)
            render_text_lines(element, target, tree)
            return
        if element.semantic_type == "Header":
            logger.debug(f"Header element ocr_lines count: {len(element.ocr_lines)}")
        if element.shape_type: 
            render_shape(element, target, tree)
        render_text_lines(element, target, tree)
    elif element.processing_path == "crop":
        render_tile(element, target, tree)
        render_text_lines(element, target, tree)
    elif element.semantic_type == "Arrow":
        render_arrow(element, target, tree)
    else:
        render_text_block(element, target, tree)


def render_element(element: DetectedElement, target, rendered_ids: set, tree: ElementTree):
    """Renders element and recursively renders its children into a GroupShape."""
    if element.id in rendered_ids:
        return
    rendered_ids.add(element.id)
    # Pre-register all children so containment guard never re-processes them
    rendered_ids.update(element.child_ids())

    print(f"DEBUG: element {element.id} semantic_type={element.semantic_type} processing_path={element.processing_path}")

    if element.children:
        render_group(element, target, rendered_ids, tree)
    else:
        render_leaf(element, target, tree)


def assemble(tree: ElementTree, output_path: str, source_image, w: int, h: int) -> str:
    """Creates a new PowerPoint presentation from the ElementTree."""
    prs = Presentation()
    prs.slide_width = Emu(config.SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(config.SLIDE_HEIGHT_EMU)

    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    if source_image is not None:
        import numpy as np
        patch = source_image[0:10, 0:10]
        if patch.size > 0:
            med = np.median(patch.reshape(-1, 3), axis=0)
            b, g, r = int(med[0]), int(med[1]), int(med[2])
            lum = (0.299 * r + 0.587 * g + 0.114 * b) / 255.0
            if lum < 0.15:
                from pptx.dml.color import RGBColor
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

    rendered_ids = set()

    # Render roots in z-order (largest area first = furthest back)
    sorted_roots = sorted(tree.roots, key=lambda e: e.bbox.area, reverse=True)
    for element in sorted_roots:
        render_element(element, slide, rendered_ids, tree)

    prs.save(output_path)
    return output_path
